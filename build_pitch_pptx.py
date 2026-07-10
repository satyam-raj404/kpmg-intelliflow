"""Build IntelliSource_Pitch.pptx — KPMG branded 60-sec pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── KPMG Brand Palette ────────────────────────────────────────────────────────
KPMG_NAVY    = RGBColor(0x00, 0x33, 0x8D)   # KPMG Blue (primary)
KPMG_MED     = RGBColor(0x00, 0x5E, 0xB8)   # Medium Blue
KPMG_LIGHT   = RGBColor(0x00, 0x91, 0xDA)   # Light Blue
KPMG_TEAL    = RGBColor(0x00, 0x99, 0xA8)   # Teal accent
KPMG_GOLD    = RGBColor(0x8F, 0x73, 0x26)   # KPMG Gold
KPMG_RED     = RGBColor(0xBC, 0x20, 0x4B)   # KPMG Red (risk/danger)
KPMG_GRAY    = RGBColor(0x63, 0x66, 0x6A)   # KPMG Gray
KPMG_LTGRAY  = RGBColor(0xF2, 0xF2, 0xF2)  # Light Gray bg
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)
GREEN_OK     = RGBColor(0x00, 0x7A, 0x33)   # Success green
ORANGE_WARN  = RGBColor(0xFF, 0x6B, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_rgb is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def kpmg_header(slide, title, subtitle=None, dark=True):
    """Standard KPMG top bar with title."""
    bg = KPMG_NAVY if dark else KPMG_LTGRAY
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), bg)
    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.08), Inches(1.15), KPMG_LIGHT)
    # Title
    add_text(slide, title,
             Inches(0.2), Inches(0.12), Inches(10), Inches(0.6),
             size=28, bold=True, color=WHITE if dark else KPMG_NAVY,
             align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.2), Inches(0.72), Inches(10), Inches(0.35),
                 size=13, bold=False, color=KPMG_LIGHT if dark else KPMG_GRAY,
                 align=PP_ALIGN.LEFT)


def kpmg_footer(slide, text="KPMG India — IntelliSource | P2P Intelligence & Analytics Platform"):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), KPMG_NAVY)
    add_text(slide, text,
             Inches(0.2), Inches(7.12), Inches(10), Inches(0.3),
             size=8, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)
    add_text(slide, "CONFIDENTIAL",
             Inches(11.5), Inches(7.12), Inches(1.6), Inches(0.3),
             size=8, bold=True, color=KPMG_GOLD, align=PP_ALIGN.RIGHT)


def divider(slide, y, color=KPMG_LIGHT):
    add_rect(slide, Inches(0.2), y, Inches(12.9), Inches(0.03), color)


def bullet_block(slide, items, x, y, w, h, size=13, color=BLACK,
                 bullet="▸", gap=Inches(0.32)):
    """Render a list of bullet strings as individual textboxes."""
    cy = y
    for item in items:
        add_text(slide, f"{bullet}  {item}", x, cy, w, Inches(0.28),
                 size=size, color=color, wrap=True)
        cy += gap
    return cy


def pill(slide, text, x, y, w, h, bg, fg=WHITE, size=11, bold=True):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, text, x, y, w, h, size=size, bold=bold,
             color=fg, align=PP_ALIGN.CENTER)


def vs_box(slide, left_title, left_items, right_title, right_items,
           y_start=Inches(1.2), left_col=None, right_col=None,
           item_color_l=KPMG_RED, item_color_r=GREEN_OK):
    """Split comparison box — SAP left, IntelliSource right."""
    lc = left_col or KPMG_GRAY
    rc = right_col or KPMG_NAVY

    col_w = Inches(6.0)
    col_h = Inches(5.6)
    gap   = Inches(0.25)
    lx    = Inches(0.2)
    rx    = lx + col_w + gap

    # Panels
    add_rect(slide, lx, y_start, col_w, col_h, KPMG_LTGRAY)
    add_rect(slide, rx, y_start, col_w, col_h, KPMG_NAVY)

    # Header bars
    add_rect(slide, lx, y_start, col_w, Inches(0.45), lc)
    add_rect(slide, rx, y_start, col_w, Inches(0.45), rc)

    add_text(slide, left_title,  lx + Inches(0.1), y_start + Inches(0.05),
             col_w - Inches(0.2), Inches(0.38),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, right_title, rx + Inches(0.1), y_start + Inches(0.05),
             col_w - Inches(0.2), Inches(0.38),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    iy = y_start + Inches(0.55)
    for item in left_items:
        add_text(slide, f"✕  {item}", lx + Inches(0.15), iy,
                 col_w - Inches(0.3), Inches(0.32),
                 size=12, color=KPMG_RED, wrap=True)
        iy += Inches(0.35)

    iy = y_start + Inches(0.55)
    for item in right_items:
        add_text(slide, f"✓  {item}", rx + Inches(0.15), iy,
                 col_w - Inches(0.3), Inches(0.32),
                 size=12, color=KPMG_LIGHT, wrap=True)
        iy += Inches(0.35)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Full navy background
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, KPMG_NAVY)

# Bold left accent stripe
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, KPMG_LIGHT)

# Diagonal accent block top-right
add_rect(s, Inches(10.5), 0, Inches(2.83), Inches(2.5), KPMG_MED)
add_rect(s, Inches(11.5), 0, Inches(1.83), Inches(1.5), KPMG_LIGHT)

# KPMG wordmark (text-based)
add_text(s, "KPMG", Inches(0.3), Inches(0.3), Inches(3), Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "India", Inches(0.3), Inches(0.95), Inches(3), Inches(0.4),
         size=16, bold=False, color=KPMG_LIGHT, align=PP_ALIGN.LEFT)

divider(s, Inches(1.5), KPMG_LIGHT)

# Main title
add_text(s, "IntelliSource", Inches(0.3), Inches(1.7), Inches(11), Inches(1.2),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s, "P2P Intelligence & Analytics Platform",
         Inches(0.3), Inches(2.85), Inches(11), Inches(0.6),
         size=22, bold=False, color=KPMG_LIGHT, align=PP_ALIGN.LEFT)

divider(s, Inches(3.6), KPMG_GOLD)

add_text(s, "Executive Demo Pitch  |  60-Second Value Proposition  |  SAP Comparison",
         Inches(0.3), Inches(3.75), Inches(11), Inches(0.4),
         size=13, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)

# Bottom bar
add_rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), KPMG_MED)
add_text(s, "KPMG India  |  Procurement Advisory  |  CONFIDENTIAL",
         Inches(0.3), Inches(6.95), Inches(10), Inches(0.45),
         size=11, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "2026",
         Inches(12.3), Inches(6.95), Inches(0.9), Inches(0.45),
         size=11, bold=True, color=KPMG_GOLD, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM: WHAT SAP CANNOT TELL YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "The Problem", "What SAP Cannot Tell You")
kpmg_footer(s)

# Intro statement
add_text(s, "SAP records every transaction. It does not tell you what went wrong.",
         Inches(0.3), Inches(1.25), Inches(12.5), Inches(0.5),
         size=16, bold=True, color=KPMG_NAVY)

divider(s, Inches(1.85), KPMG_LIGHT)

# Three pain point cards
card_y = Inches(2.0)
card_h = Inches(1.55)
cards = [
    ("🔍  No Cross-Document Visibility",
     "SOD conflicts span PO → GRN → Invoice → Payment. SAP stores each in a different module. No standard report joins them. Auditors find these conflicts — you should find them first.",
     KPMG_RED),
    ("📊  No Real-Time Governance Metrics",
     "Maverick spend, duplicate invoices, CAPEX vs. OPEX split, payment before GRN — each requires a separate transaction, a manual export, and an analyst with Excel.",
     ORANGE_WARN),
    ("⏱  Decisions Made on Stale Data",
     "By the time a weekly report is ready, the payment has cleared, the vendor is paid, and the audit finding is already written. Procurement intelligence must be live.",
     KPMG_GRAY),
]

cx = Inches(0.2)
for title, body, color in cards:
    add_rect(s, cx, card_y, Inches(4.1), card_h, KPMG_LTGRAY)
    add_rect(s, cx, card_y, Inches(4.1), Inches(0.08), color)
    add_text(s, title, cx + Inches(0.1), card_y + Inches(0.12),
             Inches(3.9), Inches(0.45), size=12, bold=True, color=color)
    add_text(s, body, cx + Inches(0.1), card_y + Inches(0.55),
             Inches(3.9), Inches(0.9), size=10.5, color=KPMG_GRAY, wrap=True)
    cx += Inches(4.3)

# The cost line
add_rect(s, Inches(0.2), Inches(3.75), Inches(12.9), Inches(0.75), KPMG_NAVY)
add_text(s,
    "One duplicate payment event: ₹5–50 lakhs. One SOD audit finding: ₹20–80 lakhs in remediation. One year of blind spots: priceless.",
    Inches(0.4), Inches(3.83), Inches(12.5), Inches(0.55),
    size=13, bold=False, color=WHITE, italic=True)

# SAP screenshot placeholder
add_rect(s, Inches(0.2), Inches(4.6), Inches(6.0), Inches(2.3), KPMG_LTGRAY)
add_rect(s, Inches(0.2), Inches(4.6), Inches(6.0), Inches(0.35), KPMG_GRAY)
add_text(s, "SAP Transaction: FBL1N — Vendor Line Items",
         Inches(0.3), Inches(4.62), Inches(5.8), Inches(0.3),
         size=9, bold=True, color=WHITE)
add_text(s,
    "[Insert SAP FBL1N screenshot — shows raw vendor line items with no anomaly flags, no duplicate detection, no cross-document context. Analyst must export and manually VLOOKUP against MB51 and ME2M to find issues.]",
    Inches(0.3), Inches(5.05), Inches(5.8), Inches(1.8),
    size=9.5, color=KPMG_GRAY, italic=True, wrap=True)

# IntelliSource screenshot placeholder
add_rect(s, Inches(6.4), Inches(4.6), Inches(6.7), Inches(2.3), KPMG_LTGRAY)
add_rect(s, Inches(6.4), Inches(4.6), Inches(6.7), Inches(0.35), KPMG_NAVY)
add_text(s, "IntelliSource — Financial Dashboard: Duplicate Invoice Detection",
         Inches(6.5), Inches(4.62), Inches(6.5), Inches(0.3),
         size=9, bold=True, color=WHITE)
add_text(s,
    "[Insert IntelliSource Financial Dashboard screenshot — shows Duplicate Invoice Rate KPI card, live count, and drill-down table with vendor, invoice ref, amount, and duplicate count — no manual work required.]",
    Inches(6.5), Inches(5.05), Inches(6.5), Inches(1.8),
    size=9.5, color=KPMG_NAVY, italic=True, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SAP vs INTELLISOURCE: SIDE-BY-SIDE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "SAP vs IntelliSource", "Where IntelliSource Fills the Gap")
kpmg_footer(s)

# Left panel — SAP
lx, rx = Inches(0.2), Inches(6.75)
pw = Inches(6.3)
py = Inches(1.25)
ph = Inches(5.5)

add_rect(s, lx, py, pw, ph, KPMG_LTGRAY)
add_rect(s, lx, py, pw, Inches(0.5), KPMG_GRAY)
add_text(s, "🔴  SAP Standard  —  What You Get Today",
         lx + Inches(0.1), py + Inches(0.07), pw - Inches(0.2), Inches(0.38),
         size=13, bold=True, color=WHITE)

# SAP screenshot placeholder block
add_rect(s, lx + Inches(0.15), py + Inches(0.6), pw - Inches(0.3), Inches(2.1), WHITE)
add_rect(s, lx + Inches(0.15), py + Inches(0.6), pw - Inches(0.3), Inches(0.28), KPMG_GRAY)
add_text(s, "SAP — ME2M (Purchase Orders) + MB51 (GRN) — Separate Transactions",
         lx + Inches(0.25), py + Inches(0.62), pw - Inches(0.5), Inches(0.25),
         size=7.5, bold=True, color=WHITE)
add_text(s,
    "[Insert side-by-side SAP ME2M and MB51 transaction screenshots]\n\nEach transaction is siloed. To answer 'Did the same person create this PO and post the GRN?' — an analyst must download both reports, match on PO number, and compare created_by fields manually. No real-time view. No alert. No audit trail.",
    lx + Inches(0.25), py + Inches(0.95), pw - Inches(0.5), Inches(1.65),
    size=9, color=KPMG_GRAY, wrap=True, italic=True)

sap_pain = [
    "6–8 transaction codes to see full P2P picture",
    "No cross-document creator comparison",
    "Maverick spend requires ME2M + VLOOKUP on PR list",
    "Duplicate invoices: no standard detection",
    "CAPEX/OPEX split: manual material group export",
    "SOD conflicts: not natively computed anywhere",
    "Data is stale — batch jobs, not live",
]
iy = py + Inches(2.85)
for item in sap_pain:
    add_text(s, f"✕  {item}", lx + Inches(0.15), iy,
             pw - Inches(0.3), Inches(0.3), size=10.5, color=KPMG_RED, wrap=True)
    iy += Inches(0.31)

# Right panel — IntelliSource
add_rect(s, rx, py, pw, ph, KPMG_NAVY)
add_rect(s, rx, py, pw, Inches(0.5), KPMG_MED)
add_text(s, "🟢  IntelliSource  —  What You Get Now",
         rx + Inches(0.1), py + Inches(0.07), pw - Inches(0.2), Inches(0.38),
         size=13, bold=True, color=WHITE)

# IntelliSource screenshot placeholder
add_rect(s, rx + Inches(0.15), py + Inches(0.6), pw - Inches(0.3), Inches(2.1), RGBColor(0x0D, 0x1F, 0x35))
add_rect(s, rx + Inches(0.15), py + Inches(0.6), pw - Inches(0.3), Inches(0.28), KPMG_MED)
add_text(s, "IntelliSource — Leadership Dashboard (SOD Conflicts + Risk Panel)",
         rx + Inches(0.25), py + Inches(0.62), pw - Inches(0.5), Inches(0.25),
         size=7.5, bold=True, color=WHITE)
add_text(s,
    "[Insert IntelliSource Leadership Dashboard screenshot]\n\nSOD Conflicts: 187 shown as live KPI. Hover → popup shows all conflicts by type (PO-Release, PO-GRN, GRN-Invoice, Invoice-Payment) with document numbers, vendor names, and users — updated the moment data is uploaded.",
    rx + Inches(0.25), py + Inches(0.95), pw - Inches(0.5), Inches(1.65),
    size=9, color=KPMG_LIGHT, wrap=True, italic=True)

is_gains = [
    "5 role-specific dashboards — one upload, all metrics live",
    "SOD conflict detection across all 4 P2P control points",
    "Maverick spend auto-flagged from PO data (no VLOOKUP)",
    "Duplicate invoice detection before payment clears",
    "CAPEX/OPEX: auto-classified, live, drillable to PO line",
    "P2P cycle time: PR → PO → GRN → Invoice → Payment",
    "Live — upload once, refresh all in < 90 seconds",
]
iy = py + Inches(2.85)
for item in is_gains:
    add_text(s, f"✓  {item}", rx + Inches(0.15), iy,
             pw - Inches(0.3), Inches(0.3), size=10.5, color=KPMG_LIGHT, wrap=True)
    iy += Inches(0.31)

# VS divider
add_rect(s, Inches(6.55), Inches(2.85), Inches(0.15), Inches(2.8), KPMG_GOLD)
add_text(s, "VS", Inches(6.47), Inches(4.1), Inches(0.35), Inches(0.4),
         size=14, bold=True, color=KPMG_GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE 5 DASHBOARDS AT A GLANCE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "Five Dashboards. Every Stakeholder Covered.", "Live KPIs from SAP Data — No Analyst Required")
kpmg_footer(s)

dashboards = [
    ("Procurement",    "Maverick Spend · PO Cycle Time · Contract Compliance · Deletion Alert",    "Procurement Manager",    KPMG_NAVY),
    ("Financial",      "Duplicate Invoices · Invoice Aging · CAPEX/OPEX Split · Overdue AP",       "Finance Controller",      KPMG_MED),
    ("Leadership",     "SOD Conflicts · High-Value POs · P2P Cycle Time · Risk Panel",             "CFO / CPO / Board",       KPMG_RED),
    ("Vendor Perf.",   "Delivery Lead Time · Top Vendor Spend · MSME Tracking · Compliance Rate",  "Vendor Management",       KPMG_TEAL),
    ("Utilization",    "CAPEX Utilization % · OPEX Utilization % · Profit Centre Drill-Down",       "Budget Owners / Finance",  KPMG_GOLD),
]

cx = Inches(0.2)
dy = Inches(1.3)
dw = Inches(2.5)
dh = Inches(5.5)
gap = Inches(0.13)

for name, kpis, audience, color in dashboards:
    add_rect(s, cx, dy, dw, dh, KPMG_LTGRAY)
    # Color top strip
    add_rect(s, cx, dy, dw, Inches(0.55), color)
    add_text(s, name, cx + Inches(0.08), dy + Inches(0.08),
             dw - Inches(0.15), Inches(0.42),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Audience pill
    add_rect(s, cx + Inches(0.08), dy + Inches(0.65), dw - Inches(0.16), Inches(0.3), color)
    add_text(s, audience, cx + Inches(0.1), dy + Inches(0.66),
             dw - Inches(0.2), Inches(0.28),
             size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # KPIs
    for i, kpi in enumerate(kpis.split(" · ")):
        ky = dy + Inches(1.05) + i * Inches(0.38)
        add_rect(s, cx + Inches(0.08), ky, dw - Inches(0.16), Inches(0.32), WHITE)
        add_text(s, kpi, cx + Inches(0.13), ky + Inches(0.03),
                 dw - Inches(0.26), Inches(0.28),
                 size=9, color=KPMG_NAVY, bold=False)
    # Screenshot placeholder
    ph_y = dy + Inches(2.7)
    add_rect(s, cx + Inches(0.08), ph_y, dw - Inches(0.16), Inches(2.55), RGBColor(0xE0, 0xE8, 0xF5))
    add_text(s, f"[Insert {name}\nDashboard\nScreenshot]",
             cx + Inches(0.08), ph_y + Inches(0.8), dw - Inches(0.16), Inches(1.0),
             size=8, color=KPMG_GRAY, italic=True, align=PP_ALIGN.CENTER)

    cx += dw + gap


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE SOD STORY (Hook)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, KPMG_NAVY)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, KPMG_RED)
kpmg_footer(s, "KPMG India — IntelliSource | SOD Conflict Detection")

add_text(s, "The Proof.", Inches(0.3), Inches(0.3), Inches(12), Inches(0.7),
         size=36, bold=True, color=WHITE)
divider(s, Inches(1.1), KPMG_GOLD)

# Stat block
add_rect(s, Inches(0.3), Inches(1.25), Inches(3.5), Inches(2.2), KPMG_RED)
add_text(s, "47", Inches(0.3), Inches(1.3), Inches(3.5), Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "SOD Conflicts\nFound — First Upload",
         Inches(0.3), Inches(2.55), Inches(3.5), Inches(0.75),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(4.0), Inches(1.25), Inches(3.0), Inches(2.2), KPMG_GRAY)
add_text(s, "11", Inches(4.0), Inches(1.3), Inches(3.0), Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Years SAP Running\nWithout Surfacing One",
         Inches(4.0), Inches(2.55), Inches(3.0), Inches(0.75),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(7.2), Inches(1.25), Inches(3.5), Inches(2.2), KPMG_MED)
add_text(s, "1", Inches(7.2), Inches(1.3), Inches(3.5), Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Day to Escalate All\n47 to Internal Audit",
         Inches(7.2), Inches(2.55), Inches(3.5), Inches(0.75),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(s, Inches(3.6), KPMG_LIGHT)

# Story text
story = (
    "One client. Large enterprise. SAP live for eleven years across 3 company codes.\n\n"
    "First upload to IntelliSource. Ninety seconds to process.\n\n"
    "Invoice vs. Payment Segregation of Duty check: same user who posted the invoice "
    "also cleared the payment — 47 times. Not detected by SAP authorisation matrix. "
    "Not caught by internal audit. Not visible in any standard report.\n\n"
    "Every conflict surfaced with the document number, vendor, user name, and date. "
    "Internal audit had all 47 escalated and under review within the same business day."
)
add_text(s, story, Inches(0.3), Inches(3.75), Inches(12.7), Inches(3.1),
         size=13, color=RGBColor(0xD0, 0xE8, 0xFF), wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE 60-SECOND PITCH (script card for presenter)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "IntelliSource in 60 Seconds", "Presenter Script")
kpmg_footer(s)

# Left: script
add_rect(s, Inches(0.2), Inches(1.25), Inches(7.8), Inches(5.55), KPMG_NAVY)
add_text(s, "PRESENTER SCRIPT  —  60 SECONDS",
         Inches(0.35), Inches(1.32), Inches(7.5), Inches(0.35),
         size=10, bold=True, color=KPMG_GOLD, align=PP_ALIGN.LEFT)

script = """\
Your Finance team runs FBL1N. Your Procurement team runs ME2M. Your audit team runs their own report. Three exports, three Excel files, three versions of the truth — and by the time someone VLOOKUPs them together, the payment is already gone.

IntelliSource eliminates that gap.

One platform. Five dashboards. Every P2P metric your CFO, CPO, and Finance Controller needs — live, from SAP data, the moment it is uploaded.

But here is what makes clients pay for it: IntelliSource finds what SAP was never designed to show you.

Your authorisation matrix prevents one person from posting and clearing a payment — in theory. IntelliSource checks every invoice and every payment in your system and surfaces each violation with a name, a document number, and a date.

One client. First upload. 47 SOD conflicts. SAP had been running eleven years. None of them visible.

That is not a dashboard. That is a governance layer your organisation does not have — and your auditors are going to ask for it.

IntelliSource. From KPMG. Because procurement intelligence should not require a data analyst.\
"""

add_text(s, script, Inches(0.35), Inches(1.75), Inches(7.55), Inches(4.85),
         size=11, color=WHITE, wrap=True)

# Right: timing guide
rx = Inches(8.2)
rw = Inches(4.9)
add_rect(s, rx, Inches(1.25), rw, Inches(5.55), KPMG_LTGRAY)
add_text(s, "TIMING GUIDE", rx + Inches(0.15), Inches(1.32),
         rw - Inches(0.3), Inches(0.35),
         size=10, bold=True, color=KPMG_NAVY)

timing = [
    ("0–10 sec",  "The problem — 3 exports, 3 versions of truth", KPMG_RED),
    ("10–20 sec", "IntelliSource intro — one platform, 5 dashboards", KPMG_MED),
    ("20–35 sec", "The differentiator — what SAP cannot show you", KPMG_NAVY),
    ("35–50 sec", "The SOD story — 47 conflicts, 11 years invisible", KPMG_GOLD),
    ("50–60 sec", "Close — governance layer, auditor demand", GREEN_OK),
]
ty = Inches(1.75)
for time, desc, color in timing:
    add_rect(s, rx + Inches(0.1), ty, rw - Inches(0.2), Inches(0.85), WHITE)
    add_rect(s, rx + Inches(0.1), ty, Inches(0.06), Inches(0.85), color)
    add_text(s, time, rx + Inches(0.25), ty + Inches(0.07),
             Inches(1.1), Inches(0.3), size=9, bold=True, color=color)
    add_text(s, desc, rx + Inches(0.25), ty + Inches(0.38),
             rw - Inches(0.5), Inches(0.42), size=9, color=KPMG_GRAY, wrap=True)
    ty += Inches(0.95)

add_text(s, "💡  Pause after the 47 number. Let it land.",
         rx + Inches(0.1), ty + Inches(0.1), rw - Inches(0.2), Inches(0.4),
         size=9, bold=True, color=KPMG_RED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATA SOURCES: FROM SAP TO INTELLISOURCE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "Data Architecture", "SAP Source → IntelliSource Table → Dashboard KPI")
kpmg_footer(s)

rows = [
    ("ME5A / EBAN",    "pr_dump",           "Purchase Requisitions",    "Procurement · Leadership"),
    ("ME2M / EKKO",    "po_dump",           "Purchase Orders",          "All 5 Dashboards"),
    ("MB51 / MKPF",    "grn_dump",          "Goods Receipts (GRN)",     "Financial · Leadership · Vendor"),
    ("MIR6 / RBKP",    "po_invoice_dump",   "PO Invoice Linkage",       "SOD Conflicts · P2P Tracker"),
    ("FBL1N / BKPF",   "invoice_dump",      "AP Invoices",              "Financial · Leadership"),
    ("F110 / BKPF",    "payment_dump",      "Outgoing Payments",        "Financial · Leadership · P2P"),
    ("XK03 / LFA1",    "vendor_master",     "Vendor Master",            "Vendor Performance"),
    ("AUT10 / CDHDR",  "change_log",        "PO Change History",        "SOD: PO Create vs Release"),
    ("ME2L / EKET",    "po_delivery_dump",  "Delivery Schedule",        "Vendor Delivery Lead Time"),
]

# Header row
hx = Inches(0.2)
hy = Inches(1.3)
col_ws = [Inches(2.5), Inches(2.7), Inches(3.5), Inches(4.3)]
col_heads = ["SAP Transaction / Table", "IntelliSource Table", "What It Captures", "Drives These KPIs"]
col_colors = [KPMG_NAVY, KPMG_MED, KPMG_GRAY, KPMG_TEAL]

cx2 = hx
for ch, cw, cc in zip(col_heads, col_ws, col_colors):
    add_rect(s, cx2, hy, cw - Inches(0.05), Inches(0.4), cc)
    add_text(s, ch, cx2 + Inches(0.08), hy + Inches(0.06),
             cw - Inches(0.15), Inches(0.3),
             size=10, bold=True, color=WHITE)
    cx2 += cw

# Data rows
for ri, (tcode, table, what, kpis) in enumerate(rows):
    ry = hy + Inches(0.42) + ri * Inches(0.52)
    bg = WHITE if ri % 2 == 0 else KPMG_LTGRAY
    cx2 = hx
    for val, cw in zip([tcode, table, what, kpis], col_ws):
        add_rect(s, cx2, ry, cw - Inches(0.05), Inches(0.48), bg)
        add_text(s, val, cx2 + Inches(0.08), ry + Inches(0.1),
                 cw - Inches(0.15), Inches(0.35),
                 size=10, color=KPMG_NAVY if val == table else BLACK,
                 bold=(val == table))
        cx2 += cw

add_text(s, "Upload CSV exports from SAP → IntelliSource ETL deduplicates, validates, classifies → all KPIs recomputed in < 90 seconds.",
         Inches(0.2), Inches(7.0), Inches(12.9), Inches(0.35),
         size=9, color=KPMG_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
kpmg_header(s, "Business Case & Next Steps", "Why Now. Why IntelliSource.")
kpmg_footer(s)

# ROI cards
roi = [
    ("₹50–250 Lakhs",  "Duplicate invoice\nrecovery (Year 1)\n0.1–0.5% of AP spend",  KPMG_RED),
    ("₹20–80 Lakhs",   "Audit finding cost\navoided per SOD\nfinding caught early",    KPMG_GOLD),
    ("800 Hours/Year", "Analyst time saved:\nmanual P2P reporting\neliminated",          KPMG_MED),
    ("3–5 Days",       "Go-live timeline\nfor single company\ncode deployment",          GREEN_OK),
]
cx3 = Inches(0.2)
for val, desc, color in roi:
    add_rect(s, cx3, Inches(1.25), Inches(3.1), Inches(2.2), KPMG_LTGRAY)
    add_rect(s, cx3, Inches(1.25), Inches(3.1), Inches(0.08), color)
    add_text(s, val, cx3 + Inches(0.1), Inches(1.38), Inches(2.9), Inches(0.75),
             size=22, bold=True, color=color, align=PP_ALIGN.LEFT)
    add_text(s, desc, cx3 + Inches(0.1), Inches(2.1), Inches(2.9), Inches(1.2),
             size=10.5, color=KPMG_GRAY, wrap=True)
    cx3 += Inches(3.3)

divider(s, Inches(3.6), KPMG_LIGHT)

# Next steps
add_text(s, "Recommended Next Steps",
         Inches(0.2), Inches(3.7), Inches(8), Inches(0.4),
         size=16, bold=True, color=KPMG_NAVY)

steps = [
    ("Week 1",  "Data readiness check — confirm SAP export access for 5 transaction codes"),
    ("Week 2",  "Pilot upload — single company code, validate KPI logic against your data"),
    ("Week 3",  "Stakeholder walkthrough — 60-minute demo with Procurement + Finance + Audit"),
    ("Week 4",  "Go/No-Go decision — full deployment or phased rollout across company codes"),
]
sy = Inches(4.2)
for wk, desc in steps:
    add_rect(s, Inches(0.2), sy, Inches(1.2), Inches(0.48), KPMG_NAVY)
    add_text(s, wk, Inches(0.2), sy + Inches(0.1), Inches(1.2), Inches(0.32),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(1.55), sy + Inches(0.1), Inches(11.2), Inches(0.32),
             size=11, color=BLACK)
    sy += Inches(0.62)

# Contact
add_rect(s, Inches(0.2), Inches(6.55), Inches(12.9), Inches(0.5), KPMG_NAVY)
add_text(s, "KPMG India — Procurement Advisory  |  getdev24@gmail.com  |  IntelliSource Pilot available immediately",
         Inches(0.35), Inches(6.62), Inches(12.5), Inches(0.35),
         size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "IntelliSource_Pitch.pptx")
prs.save(out)
print(f"Saved: {out}")
